"""PPTX report generator for SPC analysis.

Builds a 5-slide PowerPoint deck from Plotly chart figures, capability data,
and narrative text. Pure Python — no Streamlit imports.
"""

import io
import math
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Color palette (Corporate Professional) ---
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
STEEL_BLUE = RGBColor(0x4A, 0x90, 0xD9)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
ACCENT_RED = RGBColor(0xC4, 0x52, 0x3E)
MUTED_BLUE = RGBColor(0xAA, 0xBB, 0xCC)
MUTATED_GRAY = RGBColor(0x88, 0x99, 0xAA)

# --- Typography ---
FONT_FAMILY = "Calibri"


def _fig_to_image(fig, width: int = 1200, height: int = 600) -> bytes:
    """Render a Plotly Figure to PNG bytes via kaleido.

    Sanitizes HTML entities (e.g. &mdash;) in titles because kaleido's
    headless plotly.js renderer rejects them (error code 525). The browser
    handles these fine, so the visualization modules use them — but for
    headless rendering we convert to Unicode equivalents.
    """
    import copy
    import html

    fig = copy.deepcopy(fig)
    title_obj = fig.layout.title
    if title_obj and title_obj.text:
        # Decode HTML entities (&mdash; → —, &amp; → &, etc.)
        decoded = html.unescape(title_obj.text)
        if decoded != title_obj.text:
            fig.update_layout(title=dict(
                text=decoded,
                font=title_obj.font,
            ))

    return fig.to_image(format="png", width=width, height=height, scale=2)


def _add_title_shape(slide, left, top, width, height, text, font_size=36,
                     font_color=WHITE, bold=True, alignment=PP_ALIGN.LEFT):
    """Add a text box to a slide with consistent styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = FONT_FAMILY
    p.alignment = alignment
    return txBox


def _add_body_text(slide, left, top, width, height, text,
                  font_size=14, font_color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """Add a body text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.name = FONT_FAMILY
    p.alignment = alignment
    return txBox


def _add_body_paragraphs(slide, left, top, width, height, paragraphs_data):
    """Add a text box with multiple styled paragraphs.

    paragraphs_data: list of dicts with keys: text, font_size, font_color, bold, space_after
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, pd in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = pd.get("text", "")
        p.font.size = Pt(pd.get("font_size", 14))
        p.font.color.rgb = pd.get("font_color", DARK_GRAY)
        p.font.bold = pd.get("bold", False)
        p.font.name = FONT_FAMILY
        p.alignment = pd.get("alignment", PP_ALIGN.LEFT)
        if "space_after" in pd:
            p.space_after = Pt(pd["space_after"])

    return txBox


def _fill_navy_background(slide):
    """Set slide background to navy."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY


def _fill_white_background(slide):
    """Set slide background to white."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _add_kpi_table(slide, left, top, width, height, cap, n_batches):
    """Add a capability KPI table to the slide."""
    rows = 7  # header + 6 data rows
    cols = 2
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column widths
    table.columns[0].width = int(width * 0.5)
    table.columns[1].width = int(width * 0.5)

    # Header row
    for j, header in enumerate(["Metric", "Value"]):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_FAMILY
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # Data rows
    metric_rows = [
        ("Pp", f"{cap['Pp']:.2f}" if not math.isnan(cap["Pp"]) else "N/A"),
        ("Ppk", f"{cap['Ppk']:.2f}" if not math.isnan(cap["Ppk"]) else "N/A"),
        ("Mean", f"{cap['mean']:.3f}"),
        ("σ overall", f"{cap['sigma_overall']:.3f}"),
        ("PPM Total", f"{cap['total_ppm']:.0f}"),
        ("Batches", str(n_batches)),
    ]

    for i, (label, value) in enumerate(metric_rows):
        for j, text in enumerate([label, value]):
            cell = table.cell(i + 1, j)
            cell.text = text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = FONT_FAMILY
                paragraph.font.color.rgb = DARK_GRAY
                if j == 1:
                    paragraph.alignment = PP_ALIGN.RIGHT
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE

    return table_shape


def build_report(
    fig_xbar,
    fig_capability,
    fig_rolling_ppk,
    summary: dict,
    metadata: dict,
    cap: dict | None = None,
    n_batches: int = 0,
) -> io.BytesIO:
    """Build a 5-slide SPC report PPTX.

    Parameters
    ----------
    fig_xbar : plotly.graph_objects.Figure
        X-bar & R control chart figure.
    fig_capability : plotly.graph_objects.Figure
        Capability histogram figure.
    fig_rolling_ppk : plotly.graph_objects.Figure
        Rolling Ppk trend figure.
    summary : dict
        Output from reports.narrative.build_summary().
    metadata : dict
        Keys: formula, parameter, date_range (tuple), generated_at (str),
              baseline_n (int).
    cap : dict or None
        Capability metrics dict (for KPI table). If None, table is omitted.
    n_batches : int
        Number of batches included in the analysis.

    Returns
    -------
    io.BytesIO
        PPTX file bytes, ready for download.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    formula = metadata["formula"]
    parameter = metadata["parameter"]
    date_start, date_end = metadata["date_range"]
    generated_at = metadata.get("generated_at", datetime.now().strftime("%Y-%m-%d %H:%M"))
    baseline_n = metadata.get("baseline_n", 20)

    # ===== SLIDE 1: Title =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _fill_navy_background(slide1)

    _add_title_shape(
        slide1, Inches(1), Inches(2), Inches(11), Inches(1),
        "SPC Process Report", font_size=40, font_color=WHITE, alignment=PP_ALIGN.LEFT,
    )
    _add_title_shape(
        slide1, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
        f"{formula} / {parameter}", font_size=28, font_color=STEEL_BLUE,
        bold=True, alignment=PP_ALIGN.LEFT,
    )
    _add_body_text(
        slide1, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
        f"Date range: {date_start} → {date_end}", font_size=16,
        font_color=MUTED_BLUE, alignment=PP_ALIGN.LEFT,
    )
    _add_body_text(
        slide1, Inches(1), Inches(5.2), Inches(11), Inches(0.6),
        f"Generated: {generated_at}", font_size=14,
        font_color=MUTATED_GRAY, alignment=PP_ALIGN.LEFT,
    )

    # ===== SLIDE 2: Executive Summary =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_white_background(slide2)

    _add_title_shape(
        slide2, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
        "Executive Summary", font_size=28, font_color=NAVY, alignment=PP_ALIGN.LEFT,
    )

    paras = [
        {"text": summary["verdict"], "font_size": 16, "bold": True,
         "font_color": DARK_GRAY, "space_after": 12},
        {"text": summary["metrics_bullet"], "font_size": 13, "font_color": MEDIUM_GRAY,
         "space_after": 12},
        {"text": summary["trend_note"], "font_size": 13, "font_color": MEDIUM_GRAY,
         "space_after": 8},
        {"text": summary["ooc_summary"], "font_size": 13, "font_color": DARK_GRAY,
         "space_after": 8},
    ]
    if summary["action_items"]:
        paras.append({
            "text": f"Action: {summary['action_items']}",
            "font_size": 13, "font_color": ACCENT_RED,
            "bold": True, "space_after": 0,
        })

    _add_body_paragraphs(
        slide2, Inches(0.5), Inches(1.3), Inches(12), Inches(5.5), paras,
    )

    # ===== SLIDE 3: X-bar & R Chart =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_white_background(slide3)

    _add_title_shape(
        slide3, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
        "X-bar & R Control Chart", font_size=24, font_color=NAVY,
    )

    img_bytes_xbar = _fig_to_image(fig_xbar, width=1200, height=600)
    slide3.shapes.add_picture(
        io.BytesIO(img_bytes_xbar),
        Inches(0.5), Inches(0.9), Inches(12), Inches(5.5),
    )

    _add_body_text(
        slide3, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
        f"Control limits frozen from the first {baseline_n} batch(es) (baseline window).",
        font_size=10, font_color=MEDIUM_GRAY,
    )

    # ===== SLIDE 4: Capability Analysis =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_white_background(slide4)

    _add_title_shape(
        slide4, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
        "Process Capability", font_size=24, font_color=NAVY,
    )

    img_bytes_cap = _fig_to_image(fig_capability, width=800, height=500)
    slide4.shapes.add_picture(
        io.BytesIO(img_bytes_cap),
        Inches(0.3), Inches(1.0), Inches(6.5), Inches(5.5),
    )

    if cap is not None:
        _add_kpi_table(
            slide4, Inches(7.2), Inches(1.2), Inches(5.5), Inches(4.0),
            cap, n_batches,
        )

    # ===== SLIDE 5: Rolling Ppk Trend =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_white_background(slide5)

    _add_title_shape(
        slide5, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
        "Trend Analysis — Rolling Ppk", font_size=24, font_color=NAVY,
    )

    img_bytes_ppk = _fig_to_image(fig_rolling_ppk, width=1200, height=500)
    slide5.shapes.add_picture(
        io.BytesIO(img_bytes_ppk),
        Inches(0.5), Inches(0.9), Inches(12), Inches(5.5),
    )

    _add_body_text(
        slide5, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
        "Rolling Ppk (window=10) — tracks capability stability over time.",
        font_size=10, font_color=MEDIUM_GRAY,
    )

    # ===== Write to BytesIO =====
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
