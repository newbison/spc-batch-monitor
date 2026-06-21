import io
import numpy as np
import math
from pptx import Presentation
from reports.pptx_generator import build_report, _fig_to_image


def test_fig_to_image_returns_bytes():
    """Verify kaleido renders a Plotly figure to PNG bytes."""
    import plotly.graph_objects as go
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=[1, 2, 3], y=[1, 2, 3]))
    img_bytes = _fig_to_image(fig, width=800, height=400)
    assert isinstance(img_bytes, bytes)
    assert len(img_bytes) > 100  # non-trivial PNG
    # PNG magic bytes
    assert img_bytes[:4] == b'\x89PNG'


def test_fig_to_image_handles_html_entities():
    """Verify _fig_to_image sanitizes HTML entities like &mdash; that break kaleido."""
    import plotly.graph_objects as go
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=[1, 2], y=[1, 2]))
    # &mdash; causes kaleido error code 525 if not sanitized
    fig.update_layout(title="Test &mdash; Chart")
    img_bytes = _fig_to_image(fig, width=600, height=300)
    assert isinstance(img_bytes, bytes)
    assert img_bytes[:4] == b'\x89PNG'


def test_build_report_returns_bytesio():
    """Verify build_report returns a BytesIO containing a valid PPTX."""
    import plotly.graph_objects as go

    fig_xbar = go.Figure()
    fig_xbar.add_trace(go.Scatter(x=[1, 2, 3], y=[8.5, 8.4, 8.6]))
    fig_xbar.update_layout(title="X-bar & R Test", height=400)

    fig_cap = go.Figure()
    fig_cap.add_trace(go.Histogram(x=[8.5, 8.4, 8.6]))
    fig_cap.update_layout(title="Capability Test", height=300)

    fig_ppk = go.Figure()
    fig_ppk.add_trace(go.Scatter(x=[1, 2, 3], y=[1.5, 1.4, 1.6]))
    fig_ppk.update_layout(title="Rolling Ppk Test", height=300)

    summary = {
        "verdict": "Process is capable (Ppk = 1.52).",
        "metrics_bullet": "Pp = 1.80, Ppk = 1.52, Mean = 8.450, σ = 0.930, PPM = 127, N = 40 batches",
        "trend_note": "No significant trend detected.",
        "ooc_summary": "No out-of-control events detected.",
        "action_items": "",
    }
    metadata = {
        "formula": "Grade A",
        "parameter": "viscosity",
        "date_range": ("2025-01-01", "2025-03-15"),
        "generated_at": "2025-03-20 14:30",
        "baseline_n": 20,
    }

    result = build_report(
        fig_xbar=fig_xbar,
        fig_capability=fig_cap,
        fig_rolling_ppk=fig_ppk,
        summary=summary,
        metadata=metadata,
    )

    assert isinstance(result, io.BytesIO)
    result.seek(0)

    prs = Presentation(result)
    assert len(prs.slides) == 5

    # Slide 1: Title — should contain formula and parameter
    slide1_text = " ".join(
        shape.text for shape in prs.slides[0].shapes if shape.has_text_frame
    )
    assert "Grade A" in slide1_text
    assert "viscosity" in slide1_text

    # Slide 2: Summary — should contain the verdict
    slide2_text = " ".join(
        shape.text for shape in prs.slides[1].shapes if shape.has_text_frame
    )
    assert "capable" in slide2_text.lower()
    assert "1.52" in slide2_text

    # Slides 3-5: should contain images (the chart PNGs)
    for slide_idx in [2, 3, 4]:
        shapes = prs.slides[slide_idx].shapes
        has_picture = any(shape.shape_type == 13 for shape in shapes)  # MSO_SHAPE_TYPE.PICTURE
        assert has_picture, f"Slide {slide_idx + 1} should contain an embedded chart image"


def test_build_report_with_ooc_events():
    """Verify OOC events appear in summary slide text."""
    import plotly.graph_objects as go

    fig_xbar = go.Figure()
    fig_xbar.add_trace(go.Scatter(x=[1, 2], y=[8.5, 8.4]))
    fig_xbar.update_layout(height=400)

    fig_cap = go.Figure()
    fig_cap.update_layout(height=300)

    fig_ppk = go.Figure()
    fig_ppk.update_layout(height=300)

    summary = {
        "verdict": "Process is marginal (Ppk = 1.08).",
        "metrics_bullet": "Pp = 1.10, Ppk = 1.08, Mean = 8.450, PPM = 5500",
        "trend_note": "Downward trend detected in recent batches.",
        "ooc_summary": "3 out-of-control event(s) detected: B-012 (2025-03-12) through B-018 (2025-03-18).",
        "action_items": "Investigate root cause of out-of-control signals.",
    }
    metadata = {
        "formula": "Grade A", "parameter": "viscosity",
        "date_range": ("2025-03-01", "2025-03-28"),
        "generated_at": "2025-03-20", "baseline_n": 10,
    }

    result = build_report(
        fig_xbar=fig_xbar, fig_capability=fig_cap, fig_rolling_ppk=fig_ppk,
        summary=summary, metadata=metadata,
    )
    result.seek(0)
    prs = Presentation(result)

    slide2_text = " ".join(
        shape.text for shape in prs.slides[1].shapes if shape.has_text_frame
    )
    assert "marginal" in slide2_text.lower()
    assert "3 out-of-control" in slide2_text.lower()
    assert "investigate" in slide2_text.lower()
